from pptx import Presentation


def extract_images_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    images = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_image:
                image = shape.image
                # 获取图片的二进制数据
                image_bytes = image.blob
                # 将图片数据添加到列表中
                images.append(image_bytes)

    return images


if __name__ == '__main__':
    path0 = r'C:\Users\liuch\Desktop\项目开发组_周报_W01_20240606.pptx'
    print(extract_images_from_ppt(path0))
